"""
Helper functions to interact with the Presentation.

In particular, contains functions to remove, move and create various elements.
"""

import logging
import os
import sys

import copy

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.table import _Column, _Row

from .copy import textframe_copy


def export(pres, outpath):
    """
    Saves the given presentation at the given path, while deleting the previous file.
    """

    try:
        os.remove(outpath)
        logging.getLogger(__name__).info("Former %s deleted.", outpath)
    except OSError:
        pass
    pres.save(outpath)


def remove_slide(pres, idx):
    """
    Removes the slide at the given index from the presentation.
    """

    rId = pres.slides._sldIdLst[idx].rId
    pres.part.drop_rel(rId)
    del pres.slides._sldIdLst[idx]


def delete_run(r):
    """
    Deletes the given run from its paragraph.
    """
    r._r.getparent().remove(r._r)


def clear_slides(pres):
    """
    Deletes all slides from the presentation, in reverse order.
    """
    for idx in range(len(pres.slides) - 1, -1, -1):
        remove_slide(pres, idx)


def insert_slide(pres, layout, idx):
    """
    Creates a new slide with the given layout (either an actual Layout,
    or the name of a layout), and inserts it at index idx.
    """

    slide = append_slide(pres, layout)
    move_slide(pres.slides, slide, idx)
    return slide


def move_slide(slides, slide, new_idx):
    """
    Moves the given slide at the given index.
    """
    slides._sldIdLst.insert(new_idx, slides._sldIdLst[slides.index(slide)])


def fill_placeholders(slide):
    """
    Copy the content of the placeholders (actual PowerPoint placeholders,
    not placeholders as usually defined in this library) of the layout of
    the slide into the placeholders of the slide.
    """

    layout = slide.slide_layout
    for ph in layout.placeholders:
        if ph.has_text_frame:
            try:
                phid = ph.placeholder_format.idx
                tf = slide.placeholders[phid].text_frame
            except KeyError:
                continue
            textframe_copy(ph.text_frame, tf)


def append_slide(pres, layout):
    """
    Adds a slide with the given layout (either an actual layout or the name
    of an existing layout) at the end of the presentation.
    """

    if isinstance(layout, str):
        nl = pres.slide_masters[0].slide_layouts.get_by_name(layout)
        if nl is None:
            logging.getLogger(__name__).error(
                "append_slide: layout \'%s\' inexistant", layout)
            sys.exit(1)
        layout = nl

    slide = pres.slides.add_slide(layout)
    fill_placeholders(slide)
    return slide


def etalonnage(prespath, outpath):
    """
    Loads the presentation at the given path,
    then for each shape write the shape index inside,
    and lastly export the produced presentation.
    """
    pres = Presentation(prespath)
    for idx, slide in enumerate(pres.slides):
        logging.getLogger(__name__).info("Slide %d has slide_id %s", idx,
                                         slide.slide_id)

    clear_slides(pres)
    for layout in pres.slide_masters[0].slide_layouts:
        slide = pres.slides.add_slide(layout)
        for ph in slide.placeholders:
            ph.text = f"placeholder {ph.placeholder_format.idx}"
    export(pres, outpath)


def chart_replace_data(chart, values):
    """
    Updates the data in the chart.
    """

    chart_data = CategoryChartData()
    chart_data.categories = chart.plots[0].categories
    chart_data.add_series('Values', values)
    chart.replace_data(chart_data)


def find_shape(slide, name):
    """
    Finds the first shape in the slide with the given name.

    Convenient to access a shape by its name.
    """

    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def table_dup_row(table, idx, n=1, to=None):
    """
    Duplicates a table row, in particular borders and text are copied.
    - idx is the index of the row we want to copy
    - n is the number of times we want to copy the row
    - to is the index where we want to insert the row
    """

    if to is None:
        to = idx
    to = (to + len(table.rows)) % len(table.rows)

    new_row = table._tbl.tr_lst[idx]
    num = -1
    for rid, rxml in enumerate(table._tbl.iterchildren()):
        if rxml.tag == new_row.tag:
            num += 1
            if num == to:
                break

    for _ in range(n):
        new_row = copy.deepcopy(new_row)
        table._tbl.insert(rid, new_row)

    if n == 1:
        return _Row(new_row, table)
    return None


def table_dup_column(table, idx, n=1, to=None, keep_width=False):
    """
    Duplicates a column of the table, in particular borders and text are copied.
    - idx is the index of the column we want to copy
    - n is the number of times we want to copy the column
    - to is the index where we want to insert the column
    - keep_width indicates whether we want to shrink columns to add the new ones
    so that the total width is the same, or if columns are simply added.
    """

    total_width = table._graphic_frame.width
    colwidth = table.columns[idx].width
    mult = total_width / (total_width + n * colwidth)

    if to is None:
        to = idx
    to = (to + len(table.columns)) % len(table.columns)

    new_col = table._tbl.tblGrid.gridCol_lst[idx]
    num = -1
    for cid, cxml in enumerate(table._tbl.tblGrid.iterchildren()):
        if cxml.tag == new_col.tag:
            num += 1
            if num == to:
                break
    for _ in range(n):
        new_col = copy.deepcopy(new_col)
        table._tbl.tblGrid.insert(cid, new_col)

    for tr in table._tbl.tr_lst:
        new_cell = tr.tc_lst[idx]
        num = -1
        for cid, cxml in enumerate(tr.iterchildren()):
            if cxml.tag == new_cell.tag:
                num += 1
                if num == to:
                    break
        for _ in range(n):
            new_cell = copy.deepcopy(new_cell)
            tr.insert(cid, new_cell)

    if keep_width:
        for col in table.columns:
            col.width = int(col.width * mult)

    if n == 1:
        return _Column(new_col, table)
    return None
